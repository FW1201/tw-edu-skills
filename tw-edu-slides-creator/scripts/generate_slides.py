#!/usr/bin/env python3
"""教材轉精美教學簡報生成腳本 V2.0
25 種版型 × 8 種色盤 × Emoji/SVG 圖示
"""
import argparse, sys, json, urllib.request
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ══════════════════════════════════════════════════════════════
# 年級字體規格
# ══════════════════════════════════════════════════════════════
GRADE_SPECS = {
    'elementary_low':  {'title_size': 40, 'body_size': 28, 'max_bullets': 3, 'img_ratio': 0.7},
    'elementary_high': {'title_size': 36, 'body_size': 24, 'max_bullets': 4, 'img_ratio': 0.55},
    'junior':          {'title_size': 32, 'body_size': 22, 'max_bullets': 5, 'img_ratio': 0.4},
    'senior':          {'title_size': 28, 'body_size': 20, 'max_bullets': 6, 'img_ratio': 0.3},
}

# ══════════════════════════════════════════════════════════════
# 8 種色盤  {bg, primary, accent, text, light}
# ══════════════════════════════════════════════════════════════
def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

STYLES = {
    'modern':   {'bg':rgb('FFFFFF'),'primary':rgb('1A5276'),'accent':rgb('D4AC0D'),'text':rgb('1C2A35'),'light':rgb('D6EAF8')},
    'colorful': {'bg':rgb('FFFFFF'),'primary':rgb('2184C7'),'accent':rgb('E87422'),'text':rgb('2C3E50'),'light':rgb('D6EEF8')},
    'dark':     {'bg':rgb('1A1A2E'),'primary':rgb('16213E'),'accent':rgb('E94C60'),'text':rgb('EEEEEE'),'light':rgb('2D2D4A')},
    'local':    {'bg':rgb('FFFFF5'),'primary':rgb('5C3917'),'accent':rgb('C0392B'),'text':rgb('2C1A10'),'light':rgb('F5EAD5')},
    'nature':   {'bg':rgb('FFFFFF'),'primary':rgb('2E7D32'),'accent':rgb('8BC34A'),'text':rgb('1B3A1C'),'light':rgb('F1F8E9')},
    'tech':     {'bg':rgb('0D1117'),'primary':rgb('1565C0'),'accent':rgb('00BCD4'),'text':rgb('E0F2F1'),'light':rgb('1A273A')},
    'warm':     {'bg':rgb('FFFFFF'),'primary':rgb('BF360C'),'accent':rgb('FF9800'),'text':rgb('3E1C00'),'light':rgb('FFF8E1')},
    'purple':   {'bg':rgb('FFFFFF'),'primary':rgb('4A148C'),'accent':rgb('CE93D8'),'text':rgb('1A003C'),'light':rgb('F3E5F5')},
}

# ══════════════════════════════════════════════════════════════
# Emoji 圖示對照
# ══════════════════════════════════════════════════════════════
ICONS = {
    'objectives':'🎯','vocab':'📖','activity':'🛠️','timeline':'📅',
    'discussion':'💬','summary':'📝','section':'📌','image':'🖼️',
    'data':'📊','quote':'"','process':'⚙️','concept':'🧠',
    'pros':'✅','cons':'❌','agenda':'📋','closing':'🎉',
    'hero':'🌟','card':'▪','chart':'📈','funnel':'▲',
    'person':'👤','bullet':'▸','check':'✓','arrow':'→',
    'star':'★','info':'ℹ️','idea':'💡','warn':'⚠️',
}

# ══════════════════════════════════════════════════════════════
# 工具函式
# ══════════════════════════════════════════════════════════════
def get_stage(grade_str):
    g = grade_str.replace(' ', '')
    if any(x in g for x in ['一年','二年','低年']): return 'elementary_low'
    if any(x in g for x in ['三年','四年','五年','六年','中年','高年','國小']): return 'elementary_high'
    if any(x in g for x in ['國中','七','八','九']): return 'junior'
    return 'senior'

def sf(tf, text, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    """set text frame — single paragraph helper"""
    if color is None: color = rgb('1C2A35')
    tf.text = ''
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = '微軟正黑體'

def add_bg(slide, prs, color):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_rect(slide, x_cm, y_cm, w_cm, h_cm, fill_color, border_color=None):
    s = slide.shapes.add_shape(1, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color: s.line.color.rgb = border_color
    else: s.line.fill.background()
    return s

def add_tb(slide, x_cm, y_cm, w_cm, h_cm, text, size, bold=False,
           color=None, align=PP_ALIGN.LEFT, wrap=True):
    if color is None: color = rgb('1C2A35')
    tb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tb.text_frame.word_wrap = wrap
    sf(tb.text_frame, text, size, bold=bold, color=color, align=align)
    return tb

def topbar(slide, prs, color, h=1.8):
    add_rect(slide, 0, 0, prs.slide_width.cm, h, color)

def header_title(slide, prs, title, spec, c, icon=''):
    topbar(slide, prs, c['primary'])
    label = f'{icon}  {title}' if icon else title
    add_tb(slide, 0.8, 0.3, 23, 1.2, label,
           spec['title_size']-8, bold=True, color=rgb('FFFFFF'))

def bullet_block(slide, x_cm, y_cm, w_cm, h_cm, bullets, spec, c, prefix='▸  '):
    tb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tb.text_frame.word_wrap = True
    for i, b in enumerate(bullets):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run(); r.text = f'{prefix}{b}'
        r.font.size = Pt(spec['body_size']); r.font.name = '微軟正黑體'
        r.font.color.rgb = c['text']

def note_box(slide, y_cm, note, spec, c):
    add_rect(slide, 1.2, y_cm, 22.6, 1.0, rgb('F0F4F8'), c['primary'])
    add_tb(slide, 1.5, y_cm+0.1, 22, 0.8, f'💡 {note}',
           spec['body_size']-4, color=c['primary'])

# ══════════════════════════════════════════════════════════════
# 版型 1：封面
# ══════════════════════════════════════════════════════════════
def add_title_slide(prs, title, subject, grade, style_name='modern'):
    c = STYLES[style_name]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    add_rect(slide, 0, 0, prs.slide_width.cm, 2.5, c['primary'])
    add_rect(slide, 0, 16.5, prs.slide_width.cm, 2.5, c['primary'])
    for y in [2.5, 2.7]:
        add_rect(slide, 0, y, prs.slide_width.cm, 0.08, c['accent'])
    add_tb(slide, 1.5, 4, 22, 4, title, 38, bold=True,
           color=c['primary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 1.5, 8.5, 22, 2, f'{subject}｜{grade}', 20,
           color=c['text'], align=PP_ALIGN.CENTER)
    add_tb(slide, 1.5, 17.0, 22, 1.2, '108 課綱素養導向教學', 14,
           color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 版型 2：學習目標
# ══════════════════════════════════════════════════════════════
def add_objectives_slide(prs, objectives, stage, style_name='modern'):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, '本節學習目標', spec, c, '🎯')
    bullet_block(slide, 1.5, 2.5, 22, 12, objectives, spec, c, prefix='  ➤  ')

# ══════════════════════════════════════════════════════════════
# 版型 3：一般條列（content）
# ══════════════════════════════════════════════════════════════
def add_content_slide(prs, title, bullets, note='', stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    bullet_block(slide, 1.5, 2.4, 22.6, 11, bullets[:spec['max_bullets']], spec, c)
    if note: note_box(slide, 14.2, note, spec, c)

# ══════════════════════════════════════════════════════════════
# 版型 4：思考討論（discussion）
# ══════════════════════════════════════════════════════════════
def add_discussion_slide(prs, question, thinking_time='2', stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    topbar(slide, prs, rgb('1E8449'))
    add_tb(slide, 0.8, 0.3, 18, 1.2, '💬 思考時間', spec['title_size']-8,
           bold=True, color=rgb('FFFFFF'))
    add_rect(slide, 20, 0.2, 4.5, 1.2, c['accent'])
    add_tb(slide, 20.2, 0.4, 4.1, 0.8, f'⏱ {thinking_time} 分鐘', 14,
           bold=True, color=c['primary'], align=PP_ALIGN.CENTER)
    add_tb(slide, 2, 4, 21, 6, question, spec['title_size']-2,
           bold=True, color=c['text'], align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 5：重點整理（summary）
# ══════════════════════════════════════════════════════════════
def add_summary_slide(prs, title, key_points, stage='junior', style_name='modern'):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    topbar(slide, prs, rgb('D4AC0D'))
    add_tb(slide, 0.8, 0.3, 23, 1.2, f'📝 {title} — 重點整理',
           spec['title_size']-8, bold=True, color=rgb('1C2A35'))
    colors_bg = [rgb('EBF5FB'), rgb('EAF3DE'), rgb('FAEEDA'), rgb('FBEAF0')]
    for i, (key, val) in enumerate(key_points[:4]):
        col = i % 2; row = i // 2
        x = 1.2 + col * 11.9; y = 2.5 + row * 5.5
        add_rect(slide, x, y, 11.5, 5, colors_bg[i], c['primary'])
        add_tb(slide, x+0.3, y+0.3, 10.9, 1.2, key,
               spec['body_size']-2, bold=True, color=c['primary'])
        add_tb(slide, x+0.3, y+1.6, 10.9, 3, val,
               spec['body_size']-4, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 6：左右雙欄（two_column）
# ══════════════════════════════════════════════════════════════
def add_two_column_slide(prs, title, left_title='', left_bullets=None,
                         right_title='', right_text='',
                         stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    left_bullets = left_bullets or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    # 左欄
    add_rect(slide, 1.0, 2.2, 11.3, 12.5, c['light'])
    if left_title:
        add_tb(slide, 1.3, 2.4, 10.7, 1.0, left_title,
               spec['body_size']-2, bold=True, color=c['primary'])
    bullet_block(slide, 1.4, 3.6, 10.6, 11, left_bullets[:spec['max_bullets']], spec, c)
    # 右欄
    add_rect(slide, 13.1, 2.2, 11.4, 12.5, c['bg'] if c['bg'] != rgb('0D1117') else c['light'])
    if right_title:
        add_tb(slide, 13.4, 2.4, 10.7, 1.0, right_title,
               spec['body_size']-2, bold=True, color=c['accent'])
    add_tb(slide, 13.4, 3.6, 10.7, 11, right_text,
           spec['body_size'], color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 7：課堂活動（activity）
# ══════════════════════════════════════════════════════════════
def add_activity_slide(prs, title, steps=None, time_min='5',
                       grouping='小組討論', stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    steps = steps or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, f'🛠️  {title}', spec, c)
    # 步驟卡片
    step_colors = [c['primary'], c['accent'], rgb('2E7D32'), rgb('6A1B9A')]
    card_w = min(22.4 / max(len(steps), 1), 5.5)
    for i, step in enumerate(steps[:4]):
        x = 1.2 + i * (card_w + 0.3)
        add_rect(slide, x, 2.3, card_w, 10.5, step_colors[i % 4])
        add_tb(slide, x+0.2, 2.5, card_w-0.4, 1.5,
               f'Step {i+1}', spec['body_size']+2, bold=True,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.2, 4.2, card_w-0.4, 8.5, step,
               spec['body_size']-2, color=rgb('FFFFFF'),
               align=PP_ALIGN.LEFT, wrap=True)
    # 底部資訊列
    add_rect(slide, 1.0, 13.5, 22.4, 1.3, c['light'])
    add_tb(slide, 1.5, 13.6, 11, 1.0, f'⏱  時間：{time_min} 分鐘',
           spec['body_size']-2, color=c['primary'])
    add_tb(slide, 13, 13.6, 10.5, 1.0, f'👥  形式：{grouping}',
           spec['body_size']-2, color=c['primary'])

# ══════════════════════════════════════════════════════════════
# 版型 8：生字詞彙（vocab）
# ══════════════════════════════════════════════════════════════
def add_vocab_slide(prs, title, vocab_list=None, stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    vocab_list = (vocab_list or [])[:6]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, f'📖  {title}', spec, c)
    cols = 2; rows = (len(vocab_list) + 1) // 2
    cell_w = 11.3; cell_h = min(12.5 / max(rows, 1), 6.0)
    for i, v in enumerate(vocab_list):
        col = i % cols; row = i // cols
        x = 1.0 + col * (cell_w + 0.8)
        y = 2.2 + row * (cell_h + 0.3)
        add_rect(slide, x, y, cell_w, cell_h, c['light'], c['primary'])
        add_tb(slide, x+0.3, y+0.2, cell_w-0.6, cell_h*0.38,
               v.get('word', ''), spec['body_size']+4, bold=True, color=c['primary'])
        phonetic = v.get('phonetic', '')
        if phonetic:
            add_tb(slide, x+0.3, y+cell_h*0.42, cell_w-0.6, cell_h*0.18,
                   phonetic, spec['body_size']-6, color=c['accent'])
        defn = v.get('definition', '')
        if defn:
            add_tb(slide, x+0.3, y+cell_h*0.58, cell_w-0.6, cell_h*0.22,
                   defn, spec['body_size']-4, color=c['text'], wrap=True)
        ex = v.get('example', '')
        if ex:
            add_tb(slide, x+0.3, y+cell_h*0.78, cell_w-0.6, cell_h*0.18,
                   f'例：{ex}', spec['body_size']-6, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 9：時間軸（timeline）
# ══════════════════════════════════════════════════════════════
def add_timeline_slide(prs, title, events=None, stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    events = (events or [])[:6]
    n = max(len(events), 1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, f'📅  {title}', spec, c)
    # 主軸線
    axis_y = 9.5
    add_rect(slide, 1.0, axis_y-0.06, 23.0, 0.12, c['primary'])
    # 節點
    step_x = 23.0 / n
    for i, ev in enumerate(events):
        cx = 1.0 + (i + 0.5) * step_x
        # 節點圓圈
        add_rect(slide, cx-0.55, axis_y-0.55, 1.1, 1.1, c['accent'])
        add_tb(slide, cx-0.55, axis_y-0.55, 1.1, 1.1,
               str(i+1), spec['body_size']-4, bold=True,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        # 年份（奇數上、偶數下）
        if i % 2 == 0:
            add_tb(slide, cx-2.5, axis_y-3.2, 5, 1.2, ev.get('year',''),
                   spec['body_size']-2, bold=True, color=c['primary'],
                   align=PP_ALIGN.CENTER)
            add_tb(slide, cx-2.5, axis_y-2.0, 5, 1.8, ev.get('label',''),
                   spec['body_size']-4, color=c['text'],
                   align=PP_ALIGN.CENTER, wrap=True)
        else:
            add_tb(slide, cx-2.5, axis_y+1.0, 5, 1.2, ev.get('year',''),
                   spec['body_size']-2, bold=True, color=c['primary'],
                   align=PP_ALIGN.CENTER)
            add_tb(slide, cx-2.5, axis_y+2.2, 5, 1.8, ev.get('label',''),
                   spec['body_size']-4, color=c['text'],
                   align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 10：比較對照（comparison）
# ══════════════════════════════════════════════════════════════
def add_comparison_slide(prs, title, left_label='', left_items=None,
                         right_label='', right_items=None,
                         stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    left_items = left_items or []; right_items = right_items or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    # 左欄（primary 系）
    add_rect(slide, 0.8, 2.2, 11.6, 12.5, c['light'])
    add_rect(slide, 0.8, 2.2, 11.6, 1.5, c['primary'])
    add_tb(slide, 1.0, 2.3, 11.2, 1.2, left_label,
           spec['body_size'], bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
    bullet_block(slide, 1.2, 4.0, 10.8, 10, left_items[:spec['max_bullets']], spec, c)
    # 右欄（accent 系）
    accent_light = rgb('FDEBD0') if c['accent'] == rgb('D4AC0D') else c['light']
    add_rect(slide, 13.0, 2.2, 11.6, 12.5, accent_light)
    add_rect(slide, 13.0, 2.2, 11.6, 1.5, c['accent'])
    add_tb(slide, 13.2, 2.3, 11.2, 1.2, right_label,
           spec['body_size'], bold=True, color=c['primary'], align=PP_ALIGN.CENTER)
    bullet_block(slide, 13.4, 4.0, 10.8, 10, right_items[:spec['max_bullets']], spec, c)

# ══════════════════════════════════════════════════════════════
# 版型 11：大字金句（big_quote）
# ══════════════════════════════════════════════════════════════
def add_big_quote_slide(prs, quote='', source='', stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    # 大引號裝飾
    add_tb(slide, 0.5, 1.0, 5, 4, '\u201c', 96, bold=True,
           color=c['light'] if c['light'] != rgb('2D2D4A') else rgb('3A3A5A'),
           align=PP_ALIGN.LEFT)
    # 金句本文
    add_tb(slide, 2.5, 3.5, 20, 8, quote, spec['title_size']+2,
           bold=True, color=c['primary'], align=PP_ALIGN.CENTER, wrap=True)
    # 出處
    if source:
        add_rect(slide, 8, 13.5, 8.4, 1.2, c['accent'])
        add_tb(slide, 8.2, 13.6, 8, 1.0, f'—— {source}',
               spec['body_size']-2, bold=True,
               color=c['primary'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 版型 12：數據亮點（data_highlight）
# ══════════════════════════════════════════════════════════════
def add_data_highlight_slide(prs, title='', data_points=None,
                             stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    data_points = (data_points or [])[:4]
    n = max(len(data_points), 1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, f'📊  {title}', spec, c)
    card_w = 23.0 / n - 0.4
    dp_colors = [c['primary'], c['accent'], rgb('2E7D32'), rgb('6A1B9A')]
    for i, dp in enumerate(data_points):
        x = 0.8 + i * (card_w + 0.4)
        add_rect(slide, x, 2.4, card_w, 12, dp_colors[i % 4])
        icon = dp.get('icon', '📊')
        add_tb(slide, x+0.2, 3.0, card_w-0.4, 1.8, icon, 32,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.2, 4.8, card_w-0.4, 3.5, dp.get('value', ''),
               spec['title_size']+6, bold=True,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.2, 8.5, card_w-0.4, 5.5, dp.get('label', ''),
               spec['body_size']-2, color=rgb('FFFFFF'),
               align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 13：圖片聚焦（image_focus）
# ══════════════════════════════════════════════════════════════
def add_image_focus_slide(prs, title='', caption='', image_hint='',
                          stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    # 圖片佔位框
    add_rect(slide, 1.0, 2.2, 15.0, 12.5, c['light'], c['primary'])
    hint_text = image_hint if image_hint else '🖼️  請在此插入圖片'
    add_tb(slide, 3.5, 7.5, 10, 2.0, hint_text,
           spec['body_size'], color=c['primary'], align=PP_ALIGN.CENTER)
    # 右側說明欄
    add_rect(slide, 16.8, 2.2, 7.6, 12.5, c['light'])
    add_tb(slide, 17.0, 2.4, 7.2, 1.0, '說明', spec['body_size']-2,
           bold=True, color=c['primary'])
    add_tb(slide, 17.0, 3.6, 7.2, 11, caption,
           spec['body_size']-2, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 14：流程步驟（process_flow）
# ══════════════════════════════════════════════════════════════
def add_process_flow_slide(prs, title='', steps=None,
                           stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    steps = (steps or [])[:5]
    n = max(len(steps), 1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, f'⚙️  {title}', spec, c)
    box_w = (23.0 - (n-1)*1.0) / n
    step_colors = [c['primary'], c['accent'], rgb('1E8449'), rgb('7B241C'), rgb('1A5276')]
    for i, st in enumerate(steps):
        x = 0.8 + i * (box_w + 1.0)
        add_rect(slide, x, 2.4, box_w, 11.5, step_colors[i % 5])
        icon = st.get('icon', f'0{i+1}' if i < 9 else str(i+1))
        add_tb(slide, x+0.2, 2.8, box_w-0.4, 1.5, icon, 28,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.2, 4.5, box_w-0.4, 1.5,
               st.get('label', f'步驟{i+1}'), spec['body_size']-2,
               bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.2, 6.2, box_w-0.4, 7.5, st.get('desc', ''),
               spec['body_size']-4, color=rgb('FFFFFF'),
               align=PP_ALIGN.LEFT, wrap=True)
        # 箭頭
        if i < n - 1:
            add_tb(slide, x+box_w+0.1, 7.5, 0.8, 1.0, '→', 22,
                   bold=True, color=c['accent'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 版型 15：三欄分類 KWL（kanban）
# ══════════════════════════════════════════════════════════════
def add_kanban_slide(prs, title='', columns=None,
                    stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    columns = columns or [
        {'header': '已知 (K)', 'items': []},
        {'header': '想知道 (W)', 'items': []},
        {'header': '學到 (L)', 'items': []},
    ]
    columns = columns[:3]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    col_colors_bg = [c['light'], rgb('EAF3DE'), rgb('FAEEDA')]
    col_colors_hd = [c['primary'], rgb('1E8449'), rgb('D35400')]
    col_w = 7.6
    for i, col in enumerate(columns):
        x = 0.8 + i * (col_w + 0.2)
        add_rect(slide, x, 2.2, col_w, 12.5, col_colors_bg[i])
        add_rect(slide, x, 2.2, col_w, 1.6, col_colors_hd[i])
        add_tb(slide, x+0.2, 2.3, col_w-0.4, 1.3,
               col.get('header', f'欄{i+1}'), spec['body_size']-2,
               bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        items = col.get('items', [])
        for j, item in enumerate(items[:6]):
            add_rect(slide, x+0.3, 4.1+j*1.5, col_w-0.6, 1.3, rgb('FFFFFF'), c['primary'])
            add_tb(slide, x+0.5, 4.2+j*1.5, col_w-1.0, 1.1, item,
                   spec['body_size']-4, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 16：概念放射（concept_web）
# ══════════════════════════════════════════════════════════════
def add_concept_web_slide(prs, center='', branches=None,
                          stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    branches = (branches or [])[:6]
    n = len(branches)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    topbar(slide, prs, c['primary'])
    # 中心圓
    cx, cy = 12.5, 10.0
    add_rect(slide, cx-2.5, cy-1.8, 5.0, 3.6, c['primary'])
    add_tb(slide, cx-2.3, cy-1.6, 4.6, 3.2, center,
           spec['body_size'], bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER, wrap=True)
    # 分支節點
    import math
    positions = []
    for i in range(n):
        angle = math.radians(-90 + i * 360 / n)
        rx, ry = 8.5 * math.cos(angle), 5.5 * math.sin(angle)
        positions.append((cx + rx, cy + ry))
    branch_colors = [c['accent'], rgb('1E8449'), rgb('7B241C'),
                     rgb('1A5276'), rgb('6A1B9A'), rgb('BF360C')]
    for i, (bx, by) in enumerate(positions):
        add_rect(slide, bx-2.0, by-1.4, 4.0, 2.8, branch_colors[i % 6])
        label = branches[i].get('label', '') if isinstance(branches[i], dict) else str(branches[i])
        detail = branches[i].get('detail', '') if isinstance(branches[i], dict) else ''
        add_tb(slide, bx-1.8, by-1.2, 3.6, 1.2, label,
               spec['body_size']-4, bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        if detail:
            add_tb(slide, bx-1.8, by-0.1, 3.6, 1.4, detail,
                   spec['body_size']-6, color=rgb('FFFFFF'),
                   align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 17：優劣分析（pros_cons）
# ══════════════════════════════════════════════════════════════
def add_pros_cons_slide(prs, title='', pros=None, cons=None,
                        stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    pros = pros or []; cons = cons or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    # 優點欄（綠）
    add_rect(slide, 0.8, 2.2, 11.6, 12.5, rgb('E8F5E9'))
    add_rect(slide, 0.8, 2.2, 11.6, 1.8, rgb('2E7D32'))
    add_tb(slide, 1.0, 2.3, 11.2, 1.4, '✅  優點 / 支持',
           spec['body_size'], bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
    for i, p in enumerate(pros[:spec['max_bullets']]):
        add_tb(slide, 1.2, 4.3+i*1.8, 11.0, 1.6, f'✓  {p}',
               spec['body_size']-2, color=rgb('1B5E20'), wrap=True)
    # 缺點欄（紅）
    add_rect(slide, 13.0, 2.2, 11.6, 12.5, rgb('FFEBEE'))
    add_rect(slide, 13.0, 2.2, 11.6, 1.8, rgb('C62828'))
    add_tb(slide, 13.2, 2.3, 11.2, 1.4, '❌  缺點 / 反對',
           spec['body_size'], bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
    for i, co in enumerate(cons[:spec['max_bullets']]):
        add_tb(slide, 13.4, 4.3+i*1.8, 11.0, 1.6, f'✗  {co}',
               spec['body_size']-2, color=rgb('B71C1C'), wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 18：章節封面（section_cover）
# ══════════════════════════════════════════════════════════════
def add_section_cover_slide(prs, section_num='', section_title='',
                            style_name='modern', **kw):
    c = STYLES[style_name]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['primary'])
    # 裝飾條
    add_rect(slide, 0, 8.5, prs.slide_width.cm, 0.2, c['accent'])
    add_rect(slide, 0, 8.7, prs.slide_width.cm, 0.2, rgb('FFFFFF') if c['bg'] != rgb('0D1117') else c['light'])
    add_tb(slide, 2, 4.5, 20, 2.5, section_num, 28,
           color=c['accent'], align=PP_ALIGN.CENTER)
    add_tb(slide, 2, 7.0, 20, 3.5, section_title, 36,
           bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 版型 19：英雄訊息（hero_message）
# ══════════════════════════════════════════════════════════════
def add_hero_message_slide(prs, headline='', subtext='',
                           style_name='modern', stage='junior', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['primary'])
    # 橫向裝飾線
    add_rect(slide, 2, 10.5, 20.4, 0.15, c['accent'])
    add_tb(slide, 1.5, 3.5, 21.4, 5.5, headline,
           spec['title_size']+8, bold=True,
           color=rgb('FFFFFF'), align=PP_ALIGN.CENTER, wrap=True)
    if subtext:
        add_tb(slide, 2, 11.2, 20.4, 3, subtext, spec['body_size']+2,
               color=c['accent'], align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 20：三卡並列（three_cards）
# ══════════════════════════════════════════════════════════════
def add_three_cards_slide(prs, title='', cards=None,
                          stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    cards = (cards or [])[:3]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    card_colors = [c['primary'], c['accent'], rgb('2E7D32')]
    light_bg = [c['light'], rgb('FDEBD0'), rgb('E8F5E9')]
    card_w = 7.3
    for i, card in enumerate(cards):
        x = 0.8 + i * (card_w + 0.3)
        add_rect(slide, x, 2.2, card_w, 12.5, light_bg[i])
        add_rect(slide, x, 2.2, card_w, 3.5, card_colors[i])
        icon = card.get('icon', ICONS['card'])
        add_tb(slide, x+0.2, 2.3, card_w-0.4, 1.8, icon, 28,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.2, 4.2, card_w-0.4, 1.2,
               card.get('title', f'卡片 {i+1}'),
               spec['body_size']-2, bold=True, color=card_colors[i],
               align=PP_ALIGN.CENTER)
        add_tb(slide, x+0.3, 5.8, card_w-0.6, 8.5, card.get('desc', ''),
               spec['body_size']-4, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 21：四象限矩陣（four_quadrants）
# ══════════════════════════════════════════════════════════════
def add_four_quadrants_slide(prs, title='', quadrants=None,
                             stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    quadrants = (quadrants or [
        {'label': '第一象限', 'items': []},
        {'label': '第二象限', 'items': []},
        {'label': '第三象限', 'items': []},
        {'label': '第四象限', 'items': []},
    ])[:4]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    q_colors_bg = [c['light'], rgb('EAF3DE'), rgb('FAEEDA'), rgb('F4ECF7')]
    q_colors_hd = [c['primary'], rgb('1E8449'), rgb('D35400'), rgb('6A1B9A')]
    positions = [(0.8, 2.2), (12.2, 2.2), (0.8, 9.5), (12.2, 9.5)]
    for i, (x, y) in enumerate(positions):
        q = quadrants[i] if i < len(quadrants) else {'label': f'Q{i+1}', 'items': []}
        add_rect(slide, x, y, 11.0, 6.8, q_colors_bg[i])
        add_rect(slide, x, y, 11.0, 1.4, q_colors_hd[i])
        add_tb(slide, x+0.2, y+0.15, 10.6, 1.1, q.get('label', ''),
               spec['body_size']-2, bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        for j, item in enumerate(q.get('items', [])[:3]):
            add_tb(slide, x+0.4, y+1.8+j*1.6, 10.2, 1.4, f'▸  {item}',
                   spec['body_size']-4, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 22：圖示條列（icon_list）
# ══════════════════════════════════════════════════════════════
def add_icon_list_slide(prs, title='', items=None,
                        stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    items = (items or [])[:5]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    item_colors = [c['primary'], c['accent'], rgb('2E7D32'), rgb('7B241C'), rgb('6A1B9A')]
    for i, it in enumerate(items):
        y = 2.4 + i * 2.6
        add_rect(slide, 0.8, y, 2.4, 2.2, item_colors[i % 5])
        icon = it.get('icon', ICONS['bullet']) if isinstance(it, dict) else ICONS['bullet']
        add_tb(slide, 0.9, y+0.2, 2.2, 1.8, icon, 28,
               color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        heading = (it.get('heading', '') if isinstance(it, dict) else str(it))
        desc = (it.get('desc', '') if isinstance(it, dict) else '')
        add_tb(slide, 3.6, y+0.1, 20.8, 1.0, heading,
               spec['body_size'], bold=True, color=c['primary'])
        if desc:
            add_tb(slide, 3.6, y+1.1, 20.8, 1.0, desc,
                   spec['body_size']-4, color=c['text'], wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 23：數據長條（stat_bar）
# ══════════════════════════════════════════════════════════════
def add_stat_bar_slide(prs, title='', bars=None,
                       stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    bars = (bars or [])[:5]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    bar_colors = [c['primary'], c['accent'], rgb('2E7D32'), rgb('7B241C'), rgb('6A1B9A')]
    max_val = max((b.get('value', 0) for b in bars), default=100)
    bar_h = 10.5 / max(len(bars), 1)
    for i, b in enumerate(bars):
        y = 2.4 + i * (bar_h + 0.3)
        label = b.get('label', f'項目{i+1}')
        val = b.get('value', 0)
        unit = b.get('unit', '')
        add_tb(slide, 0.8, y, 6.5, bar_h, label,
               spec['body_size']-2, bold=True, color=c['text'])
        # 背景軌道
        add_rect(slide, 7.5, y+bar_h*0.2, 15.5, bar_h*0.6, c['light'], c['primary'])
        # 填充條
        fill_w = max(0.3, 15.5 * (val / max_val))
        add_rect(slide, 7.5, y+bar_h*0.2, fill_w, bar_h*0.6, bar_colors[i % 5])
        # 數值標籤
        add_tb(slide, 23.2, y, 1.4, bar_h, f'{val}{unit}',
               spec['body_size']-2, bold=True, color=c['primary'])

# ══════════════════════════════════════════════════════════════
# 版型 24：金字塔層次（pyramid）
# ══════════════════════════════════════════════════════════════
def add_pyramid_slide(prs, title='', levels=None,
                      stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    levels = (levels or [])[:5]
    n = max(len(levels), 1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, title, spec, c)
    # 色彩由深到淺（頂到底）
    alpha_list = [1.0, 0.82, 0.64, 0.48, 0.35]
    def blend(color, alpha):
        bg = c['bg']
        r = int(color.rgb[0]*alpha + bg.rgb[0]*(1-alpha))
        g = int(color.rgb[1]*alpha + bg.rgb[1]*(1-alpha))
        b2 = int(color.rgb[2]*alpha + bg.rgb[2]*(1-alpha))
        return RGBColor(min(r,255), min(g,255), min(b2,255))
    level_h = 11.0 / n
    for i, lv in enumerate(levels):
        w = 4.0 + i * (18.0 / n)
        x = (24.4 - w) / 2
        y = 2.3 + i * level_h
        fill_color = blend(c['primary'], alpha_list[min(i, 4)])
        add_rect(slide, x, y, w, level_h-0.15, fill_color)
        label = lv.get('label', '') if isinstance(lv, dict) else str(lv)
        desc = lv.get('desc', '') if isinstance(lv, dict) else ''
        text = f'{label}  {desc}' if desc else label
        add_tb(slide, x+0.3, y+0.1, w-0.6, level_h-0.3, text,
               spec['body_size']-4, bold=(i==0), color=rgb('FFFFFF'),
               align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# 版型 25：議程目錄（agenda）
# ══════════════════════════════════════════════════════════════
def add_agenda_slide(prs, title='本節課程大綱', items=None,
                     stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    items = (items or [])[:6]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    header_title(slide, prs, f'📋  {title}', spec, c)
    item_h = 12.0 / max(len(items), 1)
    for i, it in enumerate(items):
        y = 2.3 + i * (item_h + 0.2)
        num = it.get('num', f'0{i+1}') if isinstance(it, dict) else f'0{i+1}'
        label = it.get('title', str(it)) if isinstance(it, dict) else str(it)
        time_s = it.get('time', '') if isinstance(it, dict) else ''
        add_rect(slide, 0.8, y, 2.2, item_h, c['primary'])
        add_tb(slide, 0.9, y+0.1, 2.0, item_h-0.2, num,
               spec['body_size'], bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
        add_rect(slide, 3.2, y, 17.5, item_h,
                 c['light'] if i % 2 == 0 else c['bg'])
        add_tb(slide, 3.5, y+0.1, 17.0, item_h-0.2, label,
               spec['body_size'], color=c['text'])
        if time_s:
            add_rect(slide, 21.0, y, 3.4, item_h, c['accent'])
            add_tb(slide, 21.1, y+0.1, 3.2, item_h-0.2, time_s,
                   spec['body_size']-4, color=c['primary'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 版型 26：人物引言（testimonial）
# ══════════════════════════════════════════════════════════════
def add_testimonial_slide(prs, quote='', person='', role='', context='',
                          stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    topbar(slide, prs, c['primary'])
    # 大引號
    add_tb(slide, 0.8, 0.8, 4, 3, '\u201c', 80, bold=True,
           color=c['light'], align=PP_ALIGN.LEFT)
    # 引文
    add_tb(slide, 2.5, 2.5, 20, 7.5, quote, spec['title_size']-4,
           bold=False, color=c['text'], align=PP_ALIGN.CENTER, wrap=True)
    # 人物資訊卡
    add_rect(slide, 7.5, 11.0, 9.4, 3.5, c['primary'])
    add_rect(slide, 7.5, 11.0, 2.8, 3.5, c['accent'])
    add_tb(slide, 7.6, 11.5, 2.6, 2.5, '👤', 28,
           color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
    add_tb(slide, 10.5, 11.2, 6.2, 1.2, person,
           spec['body_size'], bold=True, color=rgb('FFFFFF'))
    if role:
        add_tb(slide, 10.5, 12.4, 6.2, 0.9, role,
               spec['body_size']-4, color=c['accent'])
    if context:
        add_tb(slide, 10.5, 13.3, 6.2, 0.9, context,
               spec['body_size']-6, color=rgb('CCCCCC'))

# ══════════════════════════════════════════════════════════════
# 版型 27：非對稱分割（split_screen）
# ══════════════════════════════════════════════════════════════
def add_split_screen_slide(prs, left_text='', left_sub='', right_bullets=None,
                           stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    right_bullets = right_bullets or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    # 左半色塊
    add_rect(slide, 0, 0, 11.5, prs.slide_height.cm, c['primary'])
    add_tb(slide, 0.8, 3.5, 9.8, 5.5, left_text, spec['title_size']+2,
           bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.LEFT, wrap=True)
    if left_sub:
        add_rect(slide, 0.8, 10.5, 4, 0.12, c['accent'])
        add_tb(slide, 0.8, 11.0, 9.8, 2, left_sub, spec['body_size'],
               color=c['accent'], wrap=True)
    # 右半內容
    add_tb(slide, 12.5, 3.0, 11.5, 1.5, '重點說明',
           spec['body_size'], bold=True, color=c['primary'])
    bullet_block(slide, 12.5, 5.0, 11.0, 10, right_bullets[:spec['max_bullets']], spec, c)

# ══════════════════════════════════════════════════════════════
# 版型 28：空白模板（blank_canvas）
# ══════════════════════════════════════════════════════════════
def add_blank_canvas_slide(prs, title='', prompt='',
                           stage='junior', style_name='modern', **kw):
    c = STYLES[style_name]; spec = GRADE_SPECS[stage]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['bg'])
    topbar(slide, prs, c['primary'])
    add_tb(slide, 0.8, 0.3, 23, 1.2, title,
           spec['title_size']-8, bold=True, color=rgb('FFFFFF'))
    if prompt:
        add_tb(slide, 1.0, 2.2, 22.4, 1.0, prompt,
               spec['body_size']-4, color=c['accent'])
    # 虛線框
    add_rect(slide, 1.0, 3.5, 22.4, 11.2, c['light'], c['primary'])
    add_tb(slide, 8, 8, 8.4, 2.5, '（師生共同填寫）', spec['body_size'],
           color=c['primary'], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 版型結尾頁（closing）
# ══════════════════════════════════════════════════════════════
def add_closing_slide(prs, homework='', style_name='modern'):
    c = STYLES[style_name]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, c['primary'])
    add_rect(slide, 0, 8.5, prs.slide_width.cm, 0.15, c['accent'])
    add_tb(slide, 1.5, 4.5, 21.4, 2.5, '課後任務 🎉', 30,
           bold=True, color=rgb('FFFFFF'), align=PP_ALIGN.CENTER)
    add_tb(slide, 1.5, 8.0, 21.4, 4, homework, 18,
           color=c['accent'], align=PP_ALIGN.CENTER, wrap=True)

# ══════════════════════════════════════════════════════════════
# Dispatch Table（25 種版型）
# ══════════════════════════════════════════════════════════════
LAYOUT_MAP = {
    'content':        add_content_slide,
    'discussion':     add_discussion_slide,
    'two_column':     add_two_column_slide,
    'activity':       add_activity_slide,
    'vocab':          add_vocab_slide,
    'timeline':       add_timeline_slide,
    'comparison':     add_comparison_slide,
    'big_quote':      add_big_quote_slide,
    'data_highlight': add_data_highlight_slide,
    'image_focus':    add_image_focus_slide,
    'process_flow':   add_process_flow_slide,
    'kanban':         add_kanban_slide,
    'concept_web':    add_concept_web_slide,
    'pros_cons':      add_pros_cons_slide,
    'section_cover':  add_section_cover_slide,
    'hero_message':   add_hero_message_slide,
    'three_cards':    add_three_cards_slide,
    'four_quadrants': add_four_quadrants_slide,
    'icon_list':      add_icon_list_slide,
    'stat_bar':       add_stat_bar_slide,
    'pyramid':        add_pyramid_slide,
    'agenda':         add_agenda_slide,
    'testimonial':    add_testimonial_slide,
    'split_screen':   add_split_screen_slide,
    'blank_canvas':   add_blank_canvas_slide,
}

# ══════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════
def main():
    parser = argparse.ArgumentParser(description='教材轉精美教學簡報 V2.0')
    parser.add_argument('--title',        default='課文標題')
    parser.add_argument('--subject',      default='國語文')
    parser.add_argument('--grade',        default='國中八年級')
    parser.add_argument('--stage',        default='',
                        choices=['elementary_low','elementary_high','junior','senior',''])
    parser.add_argument('--style',        default='modern',
                        choices=list(STYLES.keys()))
    parser.add_argument('--icons',        default='emoji', choices=['emoji','svg','none'])
    parser.add_argument('--content_file', default='')
    parser.add_argument('--slides_count', type=int, default=0)
    parser.add_argument('--output',       default='教學簡報.pptx')
    args = parser.parse_args()

    stage = args.stage if args.stage else get_stage(args.grade)

    # 讀取內容 JSON
    content_data = {}
    if args.content_file and Path(args.content_file).exists():
        try:
            content_data = json.loads(Path(args.content_file).read_text(encoding='utf-8'))
        except Exception as e:
            print(f'⚠️  content_file 解析失敗：{e}，使用預設內容', file=sys.stderr)

    prs = Presentation()
    prs.slide_width  = Cm(25.4)
    prs.slide_height = Cm(19.05)

    style = args.style

    # 封面
    add_title_slide(prs, args.title, args.subject, args.grade, style)

    # 學習目標
    objectives = content_data.get('objectives', [
        f'能理解《{args.title}》的主要內容與結構',
        f'能分析並解釋文章的重要概念',
        f'能運用所學內容進行討論或創作',
    ])
    add_objectives_slide(prs, objectives, stage, style)

    # 主要內容投影片（由 content_data.slides 驅動）
    slides_data = content_data.get('slides', [
        {'layout': 'hero_message', 'headline': f'今天，我們來探索《{args.title}》',
         'subtext': '準備好了嗎？讓我們一起開始！'},
        {'layout': 'content', 'title': '背景知識',
         'bullets': ['相關背景一', '相關背景二', '先備知識檢核'],
         'note': '可請學生分享先備知識'},
        {'layout': 'content', 'title': '核心概念一',
         'bullets': ['說明A', '說明B', '補充說明C'],
         'note': '使用板書或概念圖輔助'},
        {'layout': 'two_column', 'title': '概念深化',
         'left_title': '重點分析', 'left_bullets': ['面向一', '面向二'],
         'right_title': '延伸說明', 'right_text': '進一步的補充說明...'},
        {'layout': 'concept_web', 'center': '核心主題',
         'branches': [{'label': '面向一', 'detail': '說明'}, {'label': '面向二', 'detail': '說明'},
                      {'label': '面向三', 'detail': '說明'}, {'label': '面向四', 'detail': '說明'}]},
    ])

    for slide_data in slides_data:
        layout = slide_data.get('layout', 'content')
        fn = LAYOUT_MAP.get(layout, add_content_slide)
        try:
            fn(prs, stage=stage, style_name=style, **{k: v for k, v in slide_data.items() if k != 'layout'})
        except Exception as e:
            print(f'⚠️  版型 {layout} 生成失敗：{e}，改用 content', file=sys.stderr)
            add_content_slide(prs,
                title=slide_data.get('title', ''),
                bullets=slide_data.get('bullets', [str(e)]),
                stage=stage, style_name=style)

    # 討論投影片
    discussion_q = content_data.get('discussion_question',
        f'學了《{args.title}》之後，你有什麼新的發現或想法？')
    add_discussion_slide(prs, discussion_q, stage=stage, style_name=style)

    # 重點整理
    summary_points = content_data.get('summary', [
        ('核心主題', f'《{args.title}》的中心思想'),
        ('重要概念', '本課的 3 個關鍵詞'),
        ('學習技能', '閱讀理解 + 批判思考'),
        ('延伸思考', '課後可繼續探索的問題'),
    ])
    add_summary_slide(prs, args.title, summary_points[:4], stage, style)

    # 結尾
    homework = content_data.get('homework', f'請完成《{args.title}》相關學習單，下節課帶來討論。')
    add_closing_slide(prs, homework, style)

    # 儲存
    output_path = args.output
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f'✓ 教學簡報已儲存：{output_path}（{len(prs.slides)} 張投影片，版型多樣性：{len(set(s.get("layout","content") for s in slides_data))} 種）')
    print(f'  色盤：{args.style}  ｜  年級：{stage}  ｜  圖示：{args.icons}')

if __name__ == '__main__':
    main()
