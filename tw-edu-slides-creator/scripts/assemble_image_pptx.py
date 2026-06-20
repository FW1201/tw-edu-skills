#!/usr/bin/env python3
"""Assemble origin_image/slide_XX images into a PPTX with optional speaker notes."""

from __future__ import annotations

import argparse
import os
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple


def slide_images(deck_dir: Path) -> List[Path]:
    origin = deck_dir / "origin_image"
    if not origin.exists():
        return []
    pattern = re.compile(r"^slide_(\d+)\.(png|jpe?g|bmp|gif)$", re.IGNORECASE)
    found: List[Tuple[int, Path]] = []
    for path in origin.iterdir():
        if not path.is_file():
            continue
        match = pattern.match(path.name)
        if match:
            found.append((int(match.group(1)), path))
    return [path for _, path in sorted(found)]


def load_speaker_notes(deck_dir: Path) -> Dict[int, str]:
    speech = deck_dir / "speech.md"
    if not speech.exists():
        return {}
    content = speech.read_text(encoding="utf-8")
    heading = re.compile(r"^#{2,4}\s*(?:Slide\s*(\d+)|第\s*(\d+)\s*頁|第\s*(\d+)\s*页)\b.*$", re.IGNORECASE)
    notes: Dict[int, str] = {}
    current: Optional[int] = None
    lines: List[str] = []

    def flush() -> None:
        if current is None:
            return
        text = "\n".join(lines).strip()
        if text:
            notes[current] = text

    for line in content.splitlines():
        match = heading.match(line.strip())
        if match:
            flush()
            current = int(next(group for group in match.groups() if group))
            lines = []
            continue
        if current is not None:
            lines.append(line)
    flush()
    return notes


def create_pptx(images: List[Path], output: Path, aspect_ratio: str, notes: Dict[int, str]) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise SystemExit("Missing python-pptx. Install dependencies before assembling PPTX.") from exc

    prs = Presentation()
    if aspect_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    blank = prs.slide_layouts[6]
    for index, image in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
        note = notes.get(index)
        if note:
            frame = slide.notes_slide.notes_text_frame
            frame.clear()
            frame.text = note
        print(f"added slide {index}: {image.name}")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"saved {output}")
    print(f"slides={len(images)} notes={sum(1 for i in range(1, len(images) + 1) if notes.get(i))}")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("base_dir", help="Parent directory containing the deck folder.")
    parser.add_argument("output", help="Output filename, usually deck_name.pptx.")
    parser.add_argument("--aspect-ratio", choices=["16:9", "4:3"], default="16:9")
    parser.add_argument("--init", action="store_true", help="Create deck/origin_image and exit.")
    args = parser.parse_args()

    output_name = args.output if args.output.lower().endswith(".pptx") else f"{args.output}.pptx"
    deck_name = Path(output_name).stem
    base_dir = Path(args.base_dir).expanduser().resolve()
    deck_dir = base_dir / deck_name

    if args.init:
        (deck_dir / "origin_image").mkdir(parents=True, exist_ok=True)
        (deck_dir / "prompts").mkdir(parents=True, exist_ok=True)
        print(deck_dir)
        return 0

    images = slide_images(deck_dir)
    if not images:
        print(f"No final slide images found in {deck_dir / 'origin_image'}", file=sys.stderr)
        return 1

    create_pptx(images, deck_dir / output_name, args.aspect_ratio, load_speaker_notes(deck_dir))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
