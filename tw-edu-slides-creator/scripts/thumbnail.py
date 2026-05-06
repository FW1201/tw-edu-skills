#!/usr/bin/env python3
"""PPTX 縮圖格線生成器 — 供 Claude 視覺 QA 確認投影片結構
用法：python thumbnail.py --input 簡報.pptx --output 縮圖.png --cols 4
"""
import argparse, sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from PIL import Image, ImageDraw, ImageFont
except ImportError as e:
    print(f'[錯誤] 缺少依賴：{e}。請執行 pip install python-pptx Pillow', file=sys.stderr)
    sys.exit(1)

THUMB_W = 320
THUMB_H = 180
PADDING = 8
LABEL_H = 28
BG_COLOR = (30, 30, 30)
LABEL_BG = (20, 20, 20)
LABEL_FG = (200, 200, 200)


def hex_to_rgb(rgb_color):
    """將 pptx RGBColor 轉為 (r, g, b) tuple"""
    try:
        return (rgb_color[0], rgb_color[1], rgb_color[2])
    except Exception:
        return (200, 200, 200)


def get_slide_info(slide):
    """從投影片提取：背景色、頂部條色、首個文字標題"""
    bg_color = (255, 255, 255)
    bar_color = None
    title_text = ''

    # 嘗試取得背景色
    try:
        bg = slide.background.fill
        if bg.type is not None:
            bg_rgb = bg.fore_color.rgb
            bg_color = hex_to_rgb(bg_rgb)
    except Exception:
        pass

    shapes_sorted = sorted(slide.shapes, key=lambda s: (s.top or 0))
    for shape in shapes_sorted:
        # 找頂部色塊（條形）
        if bar_color is None:
            try:
                if shape.fill.type is not None and (shape.top or 0) < 2000000:
                    bar_color = hex_to_rgb(shape.fill.fore_color.rgb)
            except Exception:
                pass
        # 找第一個有文字的形狀作為標題
        if not title_text and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                title_text = text[:40]

    if bar_color is None:
        bar_color = bg_color

    return bg_color, bar_color, title_text


def draw_slide_thumb(draw, x, y, idx, bg_color, bar_color, title_text, font):
    """在 draw 物件上繪製單張縮圖色塊"""
    # 背景
    draw.rectangle([x, y, x + THUMB_W, y + THUMB_H], fill=bg_color)
    # 頂部色條（模擬 topbar）
    bar_h = max(12, THUMB_H // 8)
    draw.rectangle([x, y, x + THUMB_W, y + bar_h], fill=bar_color)
    # 頁碼標籤（左上）
    draw.rectangle([x, y, x + 28, y + 16], fill=(0, 0, 0, 160))
    draw.text((x + 3, y + 2), str(idx), fill=(255, 255, 255), font=font)
    # 標題文字（中間區域）
    if title_text:
        lines = _wrap_text(title_text, max_chars=20)
        for li, line in enumerate(lines[:3]):
            ty = y + bar_h + 8 + li * 18
            if ty + 16 < y + THUMB_H - LABEL_H:
                # 簡單陰影
                draw.text((x + 9, ty + 1), line, fill=(0, 0, 0), font=font)
                draw.text((x + 8, ty), line, fill=(230, 230, 230), font=font)
    # 底部標籤帶
    draw.rectangle([x, y + THUMB_H - LABEL_H, x + THUMB_W, y + THUMB_H], fill=LABEL_BG)
    short = title_text[:22] if title_text else f'Slide {idx}'
    draw.text((x + 4, y + THUMB_H - LABEL_H + 6), short, fill=LABEL_FG, font=font)


def _wrap_text(text, max_chars=20):
    """簡單斷行"""
    words = list(text)
    lines = []
    current = ''
    for ch in words:
        current += ch
        if len(current) >= max_chars:
            lines.append(current)
            current = ''
    if current:
        lines.append(current)
    return lines


def generate_thumbnail(input_path: str, output_path: str, cols: int = 4):
    prs = Presentation(input_path)
    n = len(prs.slides)
    if n == 0:
        print('[警告] 簡報中沒有投影片', file=sys.stderr)
        return

    rows = (n + cols - 1) // cols
    canvas_w = cols * (THUMB_W + PADDING) + PADDING
    canvas_h = rows * (THUMB_H + PADDING + LABEL_H) + PADDING + 36

    img = Image.new('RGB', (canvas_w, canvas_h), BG_COLOR)
    draw = ImageDraw.Draw(img)

    # 嘗試載入字型，失敗則用預設
    try:
        font = ImageFont.truetype('/System/Library/Fonts/PingFang.ttc', 12)
    except Exception:
        try:
            font = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 12)
        except Exception:
            font = ImageFont.load_default()

    # 標題列
    title = f'{Path(input_path).stem} — {n} 張投影片'
    draw.text((PADDING, 8), title, fill=(220, 220, 220), font=font)

    for i, slide in enumerate(prs.slides):
        col = i % cols
        row = i // cols
        x = PADDING + col * (THUMB_W + PADDING)
        y = 36 + PADDING + row * (THUMB_H + PADDING + LABEL_H)
        bg_color, bar_color, title_text = get_slide_info(slide)
        draw_slide_thumb(draw, x, y, i + 1, bg_color, bar_color, title_text, font)

    img.save(output_path)
    print(f'[完成] 縮圖格線已儲存：{output_path}（{cols} 欄 × {rows} 列，共 {n} 張）')


def main():
    parser = argparse.ArgumentParser(description='PPTX 縮圖格線生成器')
    parser.add_argument('--input',  required=True, help='輸入 .pptx 路徑')
    parser.add_argument('--output', default='thumbnail_grid.png', help='輸出 PNG 路徑')
    parser.add_argument('--cols',   type=int, default=4, help='每列張數（預設 4）')
    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f'[錯誤] 找不到檔案：{args.input}', file=sys.stderr)
        sys.exit(1)

    generate_thumbnail(args.input, args.output, args.cols)


if __name__ == '__main__':
    main()
